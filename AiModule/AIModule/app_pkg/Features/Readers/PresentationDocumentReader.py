import asyncio
import os
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation

from .BaseDocumentReader import BaseDocumentReader


class PresentationDocumentReader(BaseDocumentReader):
    """
    Reader implementation for PowerPoint presentations.

    Supports:
    - .pptx directly through python-pptx
    - .ppt by converting it to .pptx through LibreOffice
    """

    def __init__(self):
        """Initializes the presentation reader state."""
        super().__init__()
        self.slides_text: list[str] = []

    def _convert_ppt_to_pptx(self, ppt_path: str) -> str:
        """
        Converts a legacy .ppt file to .pptx using LibreOffice.

        Args:
            ppt_path: Path to the .ppt file.

        Returns:
            Path to the converted .pptx file.
        """

        temp_dir = tempfile.mkdtemp()

        subprocess.run(
            [
                "libreoffice",
                "--headless",
                "--convert-to",
                "pptx",
                "--outdir",
                temp_dir,
                ppt_path
            ],
            check=True
        )

        pptx_file = Path(temp_dir) / f"{Path(ppt_path).stem}.pptx"

        if not pptx_file.exists():
            raise RuntimeError(f"Conversion failed for {ppt_path}")

        return str(pptx_file)

    def _load_document(self) -> None:
        """
        Loads the presentation and extracts text from all slides.
        """

        if not self.file_path:
            raise ValueError("Presentation file path not configured.")

        path = Path(self.file_path)
        suffix = path.suffix.lower()

        if suffix == ".ppt":
            file_to_open = self._convert_ppt_to_pptx(self.file_path)
        elif suffix == ".pptx":
            file_to_open = self.file_path
        else:
            raise ValueError(f"Unsupported presentation format: {suffix}")

        presentation = Presentation(file_to_open)

        self.slides_text = []

        for index, slide in enumerate(presentation.slides, start=1):
            texts = []

            # Extract text from all shapes that expose a text attribute
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())

            slide_text = "\n".join(texts)

            if slide_text:
                self.slides_text.append(
                    f"=== [Slide {index}] ===\n{slide_text}"
                )

    def read_document(self) -> str:
        """
        Returns all extracted slide text.
        """

        return "\n\n".join(self.slides_text)

    def read_slides(self) -> list[str]:
        """
        Returns extracted text per slide.
        """

        return self.slides_text

    def read_pages(self, page_numbers: list[int]) -> str:
        """
        Reads selected slides.

        Args:
            page_numbers: 0-based slide indexes.

        Returns:
            Text from selected slides.
        """

        selected = []

        for num in page_numbers:
            if 0 <= num < len(self.slides_text):
                selected.append(self.slides_text[num])

        if not selected:
            raise ValueError("No valid slides to read.")

        return "\n\n".join(selected)


async def main():
    """
    Tests the PresentationDocumentReader directly.
    """

    base_dir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    testfiles_dir = os.path.join(base_dir, "Testfiles")

    test_files = [
        "example.pptx",
        "example.ppt"
    ]

    for file_name in test_files:
        file_path = os.path.join(testfiles_dir, file_name)

        if not os.path.exists(file_path):
            print(f"\nSkipping {file_name}: file does not exist.")
            continue

        reader = PresentationDocumentReader()
        reader.configure(file_path=file_path, reader_mode="document")
        print(await reader.process())

        reader = PresentationDocumentReader()
        reader.configure(file_path=file_path, reader_mode="pages", page_numbers=[0, 1])
        print(await reader.process())


if __name__ == "__main__":
    asyncio.run(main())